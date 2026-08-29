from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from statistics import median
from typing import Any

from pptx import Presentation

from slidethus.constants import SCHEMA_VERSION
from slidethus.errors import WorkflowApplicationError
from slidethus.io_utils import atomic_create_json, read_json, sha256_file, sha256_json
from slidethus.schema_registry import SchemaRegistry


@dataclass(frozen=True)
class StyleExtractionResult:
    candidate: dict[str, Any]
    path: Path
    changed: bool


def _rgb(value: Any) -> str | None:
    try:
        raw = value.rgb
    except Exception:  # noqa: BLE001
        return None
    if raw is None:
        return None
    text = str(raw).upper()
    return f"#{text}" if len(text) == 6 else None


def _neutral(color: str) -> bool:
    try:
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
    except ValueError:
        return True
    return max(r, g, b) - min(r, g, b) < 24 or max(r, g, b) > 240 or max(r, g, b) < 28


def _dark(color: str) -> bool:
    try:
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
    except ValueError:
        return False
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) < 120


def extract_pptx_style_candidate(workspace: Path, source: Path) -> StyleExtractionResult:
    """Extract reusable tokens from one PPTX without copying font or media bytes."""

    workspace = workspace.resolve()
    source = source.resolve()
    if not source.is_file() or source.suffix.lower() != ".pptx":
        raise WorkflowApplicationError("Extract Style currently requires one PPTX reference")
    try:
        presentation = Presentation(source)
    except Exception as exc:  # noqa: BLE001
        raise WorkflowApplicationError(f"Style reference PPTX cannot be opened: {exc}") from exc

    fonts: Counter[str] = Counter()
    colors: Counter[str] = Counter()
    text_colors: Counter[str] = Counter()
    sizes: list[float] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            try:
                fill_color = _rgb(shape.fill.fore_color)
            except Exception:  # noqa: BLE001
                fill_color = None
            if fill_color:
                colors[fill_color] += 1
            try:
                line_color = _rgb(shape.line.color)
            except Exception:  # noqa: BLE001
                line_color = None
            if line_color:
                colors[line_color] += 1
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    name = str(run.font.name or "").strip()
                    if name:
                        fonts[name] += 1
                    if run.font.size is not None:
                        sizes.append(float(run.font.size.pt))
                    color = _rgb(run.font.color)
                    if color:
                        colors[color] += 1
                        text_colors[color] += 1

    state = read_json(workspace / "project_state.json")
    project_id = str(state["project_id"])
    outline_path = workspace / "outline/deck_outline.json"
    if outline_path.is_file():
        deck_id = str(read_json(outline_path)["deck_id"])
    else:
        deck_id = f"DECK-{project_id}"

    font = fonts.most_common(1)[0][0] if fonts else "Aptos"
    dark_text = next((color for color, _ in text_colors.most_common() if _dark(color)), "#17233C")
    brand_colors = [color for color, _ in colors.most_common() if not _neutral(color) and color != dark_text]
    primary = brand_colors[0] if brand_colors else "#154C5A"
    accent = brand_colors[1] if len(brand_colors) > 1 else "#D76745"
    body_size = max(14.0, min(24.0, float(median(sizes)) if sizes else 20.0))
    title_size = max(28.0, min(44.0, max(sizes) if sizes else 32.0))
    display_size = max(title_size, min(56.0, title_size * 1.25))
    aspect = presentation.slide_width / max(1, presentation.slide_height)
    aspect_ratio = "16:9" if abs(aspect - 16 / 9) < 0.08 else "4:3" if abs(aspect - 4 / 3) < 0.08 else "custom"

    source_digest = sha256_file(source)
    candidate = {
        "schema_version": SCHEMA_VERSION,
        "project_id": project_id,
        "deck_id": deck_id,
        "theme_id": f"THEME-EXTRACTED-{source_digest[:12].upper()}",
        "tone": ["reference-derived", "reusable"],
        "canvas": {"background": "#FFFFFF", "aspect_ratio": aspect_ratio},
        "colors": {
            "background": "#FFFFFF",
            "surface": "#FFFFFF",
            "text_primary": dark_text,
            "text_secondary": "#667085",
            "primary": primary,
            "accent": accent,
        },
        "typography": {
            "display": {
                "font_family": font,
                "font_size": round(display_size, 1),
                "font_weight": 700,
                "line_height": 1.12,
                "color": dark_text,
                "letter_spacing": 0,
            },
            "title": {
                "font_family": font,
                "font_size": round(title_size, 1),
                "font_weight": 700,
                "line_height": 1.18,
                "color": dark_text,
                "letter_spacing": 0,
            },
            "body": {
                "font_family": font,
                "font_size": round(body_size, 1),
                "font_weight": 400,
                "line_height": 1.28,
                "color": dark_text,
                "letter_spacing": 0,
            },
            "caption": {
                "font_family": font,
                "font_size": max(10.0, round(body_size * 0.65, 1)),
                "font_weight": 400,
                "line_height": 1.2,
                "color": "#667085",
                "letter_spacing": 0,
            },
        },
        "spacing": {
            "base": 8,
            "region_gap": 24,
            "safe_area": {"top": 40, "right": 56, "bottom": 40, "left": 56},
        },
        "shape_rules": {
            "source": "pptx-reference",
            "source_filename": source.name,
            "source_sha256": source_digest,
            "rights_status": "reference_only",
            "extracted_font_count": len(fonts),
            "extracted_color_count": len(colors),
        },
        "layout_policy": {
            "max_same_family_consecutive": 3,
            "max_bento_ratio": 0.35,
            "min_gap": 20,
        },
        "forbidden_patterns": [
            "copy-unlicensed-font-bytes",
            "copy-unlicensed-brand-assets",
            "pixel-template-lock",
        ],
        "font_fallbacks": {font: ["Arial", "Helvetica", "Liberation Sans", "Noto Sans"]},
    }
    errors = list(SchemaRegistry().validator("visual_system").iter_errors(candidate))
    if errors:
        raise WorkflowApplicationError(
            "Extracted Visual System candidate is invalid: " + errors[0].message
        )
    digest = sha256_json(candidate)
    path = workspace / ".slidethus/workflows/style-candidates" / f"{digest}.json"
    changed = atomic_create_json(path, candidate)
    if not changed and read_json(path) != candidate:
        raise WorkflowApplicationError(f"Immutable style candidate conflict: {path}")
    return StyleExtractionResult(candidate=candidate, path=path, changed=changed)
