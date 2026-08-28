from __future__ import annotations

import json
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slidethus.errors import RenderBackendError
from slidethus.io_utils import (
    atomic_create_bytes,
    atomic_create_json,
    ensure_within,
    read_json,
    sha256_bytes,
)
from slidethus.protocols import RenderRequest, RenderResult
from slidethus.render_backends.node_toolchain import node_executable, validate_sidecar
from slidethus.render_backends.node_toolchain import renderer_root as resolve_renderer_root
from slidethus.services.render_assets import RenderAssetService, ResolvedRenderAsset
from slidethus.services.render_compile import RenderCompileResult, RenderCompileService

_EDITABILITY_ORDER = {"E0": 0, "E1": 1, "E2": 2, "E3": 3, "E4": 4}
_FIXED_ZIP_TIME = (1980, 1, 1, 0, 0, 0)
_CORE_NS = {
    "dcterms": "http://purl.org/dc/terms/",
}


@dataclass(frozen=True)
class PptxStructureMeasurement:
    slide_count: int
    text_shapes: int
    tables: int
    charts: int
    pictures: int
    native_shapes: int
    editability_level: str


def _normalize_core_properties(payload: bytes, generated_at: str) -> bytes:
    try:
        root = ET.fromstring(payload)
    except ET.ParseError:
        return payload
    for tag in ("created", "modified"):
        element = root.find(f"dcterms:{tag}", _CORE_NS)
        if element is not None:
            element.text = generated_at
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def normalize_pptx_archive(source: Path, *, generated_at: str) -> bytes:
    """Return deterministic PPTX bytes with stable entry order, times and core timestamps."""

    with zipfile.ZipFile(source) as archive:
        entries = sorted(archive.infolist(), key=lambda item: item.filename)
        payloads = {entry.filename: archive.read(entry.filename) for entry in entries}
    if "docProps/core.xml" in payloads:
        payloads["docProps/core.xml"] = _normalize_core_properties(
            payloads["docProps/core.xml"], generated_at
        )
    with tempfile.SpooledTemporaryFile(max_size=32 * 1024 * 1024) as handle:
        with zipfile.ZipFile(
            handle,
            mode="w",
            compression=zipfile.ZIP_DEFLATED,
            compresslevel=9,
        ) as output:
            for filename in sorted(payloads):
                info = zipfile.ZipInfo(filename=filename, date_time=_FIXED_ZIP_TIME)
                info.compress_type = zipfile.ZIP_DEFLATED
                info.create_system = 3
                info.external_attr = 0o600 << 16
                output.writestr(info, payloads[filename])
        handle.seek(0)
        return handle.read()


def measure_pptx_structure(path: Path, *, mode: str) -> PptxStructureMeasurement:
    """Reopen a real PPTX and conservatively measure its native object structure."""

    presentation = Presentation(path)
    text_shapes = 0
    tables = 0
    charts = 0
    pictures = 0
    native_shapes = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                tables += 1
            elif getattr(shape, "has_chart", False):
                charts += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            elif getattr(shape, "has_text_frame", False):
                text_shapes += 1
            else:
                native_shapes += 1
    if mode == "native" and pictures == 0:
        editability = "E3"
    else:
        editability = "E2"
    return PptxStructureMeasurement(
        slide_count=len(presentation.slides),
        text_shapes=text_shapes,
        tables=tables,
        charts=charts,
        pictures=pictures,
        native_shapes=native_shapes,
        editability_level=editability,
    )


class _PptxGenJSBackend:
    name = ""
    version = "1.0.0"
    mode = ""

    def __init__(
        self,
        *,
        renderer_root: Path | None = None,
        node: str | None = None,
        timeout_seconds: int = 180,
        compiled: RenderCompileResult | None = None,
        assets: dict[str, ResolvedRenderAsset] | None = None,
    ) -> None:
        self.renderer_root = resolve_renderer_root(renderer_root)
        self.node = node
        self.timeout_seconds = timeout_seconds
        self.compiled = compiled
        self.assets = assets

    def render(self, request: RenderRequest) -> RenderResult:
        if request.target_format != "pptx":
            raise RenderBackendError(f"{self.name} only supports target_format=pptx")
        if request.target_editability_level not in _EDITABILITY_ORDER:
            raise RenderBackendError(
                f"Unknown target editability: {request.target_editability_level}"
            )
        workspace = request.workspace.resolve()
        output_root = ensure_within(workspace, request.output_dir.resolve())
        output_root.mkdir(parents=True, exist_ok=True)
        executable = node_executable(self.node)
        script = validate_sidecar(
            self.renderer_root,
            script_name="render.mjs",
            dependencies={"pptxgenjs": "4.0.1"},
        )
        compiled = self.compiled or RenderCompileService(workspace).compile()
        resolved_assets = self.assets or RenderAssetService(workspace).resolve(
            tuple(compiled.ir.get("asset_ids", []))
        )
        assets = {
            asset_id: asset.as_sidecar_value()
            for asset_id, asset in resolved_assets.items()
        }
        with tempfile.TemporaryDirectory(prefix="slidethus-node-render-", dir=output_root) as name:
            temporary = Path(name)
            raw_pptx = temporary / f"{self.mode}.pptx"
            sidecar_report = temporary / "sidecar-report.json"
            asset_map_path = temporary / "assets.json"
            asset_map_path.write_text(
                json.dumps(assets, ensure_ascii=False, indent=2) + "\n",
                encoding="utf-8",
            )
            process = subprocess.run(
                [
                    executable,
                    str(script),
                    "--mode",
                    self.mode,
                    "--input",
                    str(compiled.path),
                    "--assets",
                    str(asset_map_path),
                    "--output",
                    str(raw_pptx),
                    "--report",
                    str(sidecar_report),
                    "--target-editability",
                    request.target_editability_level,
                ],
                cwd=self.renderer_root,
                check=False,
                capture_output=True,
                text=True,
                timeout=self.timeout_seconds,
            )
            if process.returncode != 0 or not raw_pptx.is_file() or not sidecar_report.is_file():
                detail = (process.stderr or process.stdout).strip()
                raise RenderBackendError(
                    f"{self.name} sidecar failed: {detail or process.returncode}"
                )
            sidecar = read_json(sidecar_report)
            normalized = normalize_pptx_archive(
                raw_pptx,
                generated_at=str(compiled.ir["generated_at"]),
            )
        digest = sha256_bytes(normalized)
        backend_dir = output_root / self.name
        output_path = backend_dir / f"{compiled.ir['deck_id'].lower()}-{digest[:16]}.pptx"
        created = atomic_create_bytes(output_path, normalized)
        if not created and output_path.read_bytes() != normalized:
            raise RenderBackendError(
                f"Immutable {self.name} output contains different content: {output_path}"
            )
        measurement = measure_pptx_structure(output_path, mode=self.mode)
        if measurement.slide_count != len(compiled.ir["slides"]):
            raise RenderBackendError(
                f"{self.name} slide count mismatch: {measurement.slide_count} != {len(compiled.ir['slides'])}"
            )
        expected_counts = dict(sidecar.get("object_counts", {}))
        if int(expected_counts.get("table", 0)) != measurement.tables:
            raise RenderBackendError(f"{self.name} table count differs from sidecar report")
        if int(expected_counts.get("chart", 0)) != measurement.charts:
            raise RenderBackendError(f"{self.name} chart count differs from sidecar report")
        expected_pictures = int(expected_counts.get("image", 0)) + int(
            expected_counts.get("embedded_svg", 0)
        )
        if expected_pictures != measurement.pictures:
            raise RenderBackendError(f"{self.name} picture count differs from sidecar report")
        if sidecar.get("measured_editability_level") != measurement.editability_level:
            raise RenderBackendError(
                f"{self.name} editability measurement disagrees with reopened PPTX"
            )
        report_data = {
            **sidecar,
            "output_path": output_path.relative_to(workspace).as_posix(),
            "output_sha256": digest,
            "structure_measurement": {
                "slide_count": measurement.slide_count,
                "text_shapes": measurement.text_shapes,
                "tables": measurement.tables,
                "charts": measurement.charts,
                "pictures": measurement.pictures,
                "native_shapes": measurement.native_shapes,
            },
        }
        report_path = backend_dir / f"{digest[:16]}-measurement.json"
        report_created = atomic_create_json(report_path, report_data)
        if not report_created and read_json(report_path) != report_data:
            raise RenderBackendError(
                f"Immutable {self.name} measurement report contains different content"
            )
        warnings = list(sidecar.get("warnings", []))
        if (
            _EDITABILITY_ORDER[measurement.editability_level]
            < _EDITABILITY_ORDER[request.target_editability_level]
            and not warnings
        ):
            warnings.append(
                f"{self.name} measured {measurement.editability_level}, below requested "
                f"{request.target_editability_level}."
            )
        return RenderResult(
            status="success",
            output_paths=(output_path, report_path),
            actual_editability_level=measurement.editability_level,
            warnings=tuple(str(item) for item in warnings),
        )


class PptxGenJSNativeRenderBackend(_PptxGenJSBackend):
    """Render native text, shapes, tables and charts through PptxGenJS."""

    name = "pptxgenjs-native"
    mode = "native"


class PptxGenJSHybridRenderBackend(_PptxGenJSBackend):
    """Render native text plus embedded complex SVG/image objects through PptxGenJS."""

    name = "pptxgenjs-hybrid"
    mode = "hybrid"
