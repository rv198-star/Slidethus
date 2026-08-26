from __future__ import annotations

import html
import math
import os
import re
import shutil
import subprocess
import tempfile
from collections.abc import Sequence
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from slidethus.io_utils import atomic_write_bytes, atomic_write_json, ensure_within, read_json
from slidethus.protocols import RenderRequest, RenderResult

LOGICAL_WIDTH = 1280
LOGICAL_HEIGHT = 720
SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5


def _hex_color(value: str) -> RGBColor:
    normalized = value.lstrip("#")
    if len(normalized) != 6:
        raise ValueError(f"Expected six-digit color, got {value!r}")
    return RGBColor.from_string(normalized.upper())


def _inches(value: float, logical_total: int, physical_total: float) -> Inches:
    return Inches(value / logical_total * physical_total)


def _block_text(content: Any) -> list[str]:
    if isinstance(content, list):
        return [str(item) for item in content]
    if isinstance(content, dict):
        return [f"{key}: {value}" for key, value in content.items()]
    return [str(content)]


def _xml_escape(value: str) -> str:
    return html.escape(value, quote=True)


def _set_run_font(run: Any, style: dict[str, Any]) -> None:
    """Set Latin and East Asian typefaces for cross-renderer CJK fidelity."""

    run.font.name = style["font_family"]
    run.font.size = Pt(style["font_size"])
    run.font.bold = style["font_weight"] >= 600
    run.font.color.rgb = _hex_color(style["color"])
    properties = run._r.get_or_add_rPr()  # python-pptx has no public East Asian API
    east_asian = properties.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if east_asian is None:
        east_asian = OxmlElement("a:ea")
        properties.append(east_asian)
    east_asian.set("typeface", style["font_family"])


def _style_for(block: dict[str, Any], visual: dict[str, Any]) -> dict[str, Any]:
    if block["semantic_role"] == "headline":
        return visual["typography"][
            "display" if block["block_id"].startswith("BLK-S001-") else "title"
        ]
    if block["semantic_role"] in {"caption", "footer", "label"}:
        return visual["typography"]["caption"]
    return visual["typography"]["body"]


def _rectangles_overlap(left: dict[str, Any], right: dict[str, Any]) -> bool:
    return not (
        left["x"] + left["w"] <= right["x"]
        or right["x"] + right["w"] <= left["x"]
        or left["y"] + left["h"] <= right["y"]
        or right["y"] + right["h"] <= left["y"]
    )


def build_layout_diagnostics(workspace: Path, output_path: Path) -> dict[str, Any]:
    """Measure planning geometry and estimated text capacity before final design."""

    workspace = workspace.resolve()
    specs = read_json(workspace / "slides/slide_specs.json")
    layouts = read_json(workspace / "layout/layout_plans.json")
    visual = read_json(workspace / "design/visual_system.json")
    specs_by_id = {slide["slide_id"]: slide for slide in specs["slides"]}
    width = layouts["canvas"]["width"]
    height = layouts["canvas"]["height"]
    safe = layouts["safe_area"]
    slide_reports: list[dict[str, Any]] = []
    all_issues: list[dict[str, Any]] = []

    for plan in layouts["plans"]:
        slide = specs_by_id[plan["slide_id"]]
        blocks = {block["block_id"]: block for block in slide["content_blocks"]}
        checks: list[dict[str, Any]] = []
        for region in plan["regions"]:
            block = blocks[region["block_id"]]
            style = _style_for(block, visual)
            values = _block_text(block["content"])
            char_count = sum(len(value) for value in values)
            usable_width = max(1, region["w"] - 32)
            usable_height = max(1, region["h"] - 28)
            chars_per_line = max(1, math.floor(usable_width / (style["font_size"] * 0.78)))
            lines = max(len(values), math.ceil(char_count / chars_per_line))
            line_capacity = max(
                1, math.floor(usable_height / (style["font_size"] * style["line_height"]))
            )
            font_floor_applies = block["semantic_role"] not in {"caption", "footer", "label"}
            meets_font_floor = (
                not font_floor_applies
                or style["font_size"] >= slide["density_budget"]["min_body_pt"]
            )
            within_canvas = (
                region["x"] >= 0
                and region["y"] >= 0
                and region["x"] + region["w"] <= width
                and region["y"] + region["h"] <= height
            )
            within_safe_area = (
                region["x"] >= safe["left"]
                and region["y"] >= safe["top"]
                and region["x"] + region["w"] <= width - safe["right"]
                and region["y"] + region["h"] <= height - safe["bottom"]
            )
            fits_estimate = lines <= line_capacity
            check = {
                "region_id": region["region_id"],
                "block_id": region["block_id"],
                "within_canvas": within_canvas,
                "within_safe_area": within_safe_area,
                "estimated_lines": lines,
                "line_capacity": line_capacity,
                "font_size_pt": style["font_size"],
                "min_body_pt": slide["density_budget"]["min_body_pt"],
                "meets_font_floor": meets_font_floor,
                "fits_estimate": fits_estimate,
            }
            checks.append(check)
            if not within_canvas or not within_safe_area or not fits_estimate or not meets_font_floor:
                all_issues.append(
                    {
                        "slide_id": plan["slide_id"],
                        "region_id": region["region_id"],
                        "kind": (
                            "bounds"
                            if not within_canvas or not within_safe_area
                            else "font_floor"
                            if not meets_font_floor
                            else "overflow"
                        ),
                        "severity": "major",
                    }
                )
        for index, left in enumerate(plan["regions"]):
            for right in plan["regions"][index + 1 :]:
                if _rectangles_overlap(left, right):
                    all_issues.append(
                        {
                            "slide_id": plan["slide_id"],
                            "region_id": left["region_id"],
                            "other_region_id": right["region_id"],
                            "kind": "collision",
                            "severity": "major",
                        }
                    )
        slide_reports.append(
            {
                "slide_id": plan["slide_id"],
                "layout_family": plan["layout_family"],
                "checks": checks,
            }
        )

    report = {
        "schema_version": "0.1.0",
        "kind": "layout_diagnostics",
        "project_id": specs["project_id"],
        "canvas": layouts["canvas"],
        "safe_area": safe,
        "slides": slide_reports,
        "issues": all_issues,
        "status": "pass" if not all_issues else "fail",
    }
    atomic_write_json(output_path, report)
    return report


def _new_presentation() -> tuple[Presentation, Any]:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    return presentation, presentation.slide_layouts[6]


class DebugPptxRenderBackend:
    """Compile Layout Plans into an inspectable PPTX with IDs and diagnostics."""

    name = "python-pptx-layout-debug"
    version = "0.2.0"

    def render(self, request: RenderRequest) -> RenderResult:
        if request.target_format != "pptx":
            raise ValueError("DebugPptxRenderBackend only supports pptx")
        workspace = request.workspace.resolve()
        output_dir = ensure_within(workspace, request.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        specs = read_json(workspace / "slides/slide_specs.json")
        layouts = read_json(workspace / "layout/layout_plans.json")
        visual = read_json(workspace / "design/visual_system.json")
        project_id = specs["project_id"]
        output_path = output_dir / f"{project_id.lower()}-debug.pptx"
        specs_by_id = {slide["slide_id"]: slide for slide in specs["slides"]}
        presentation, blank_layout = _new_presentation()

        for plan in layouts["plans"]:
            slide = presentation.slides.add_slide(blank_layout)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = _hex_color("#F3F4F6")
            self._add_grid(slide)
            self._add_safe_area(slide, layouts["safe_area"])
            blocks = {
                block["block_id"]: block
                for block in specs_by_id[plan["slide_id"]]["content_blocks"]
            }
            for index, region in enumerate(plan["regions"], start=1):
                self._add_debug_region(
                    slide,
                    region,
                    blocks[region["block_id"]],
                    index,
                    visual["typography"]["body"]["font_family"],
                )
            self._add_debug_label(slide, plan)

        buffer = BytesIO()
        presentation.save(buffer)
        atomic_write_bytes(output_path, buffer.getvalue())
        self.validate_output(output_path, layouts)
        return RenderResult(
            status="success",
            output_paths=(output_path,),
            actual_editability_level="E3",
            warnings=("Debug PPTX is a diagnostic artifact, not a designed deliverable.",),
        )

    @staticmethod
    def _add_grid(slide: Any) -> None:
        for x in range(64, LOGICAL_WIDTH, 64):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                _inches(x, LOGICAL_WIDTH, SLIDE_WIDTH_INCHES),
                Inches(0),
                Inches(0.005),
                Inches(SLIDE_HEIGHT_INCHES),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = _hex_color("#E5E7EB")
            line.line.fill.background()
        for y in range(64, LOGICAL_HEIGHT, 64):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                _inches(y, LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES),
                Inches(SLIDE_WIDTH_INCHES),
                Inches(0.005),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = _hex_color("#E5E7EB")
            line.line.fill.background()

    @staticmethod
    def _add_safe_area(slide: Any, safe: dict[str, Any]) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            _inches(safe["left"], LOGICAL_WIDTH, SLIDE_WIDTH_INCHES),
            _inches(safe["top"], LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES),
            _inches(LOGICAL_WIDTH - safe["left"] - safe["right"], LOGICAL_WIDTH, SLIDE_WIDTH_INCHES),
            _inches(LOGICAL_HEIGHT - safe["top"] - safe["bottom"], LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES),
        )
        shape.fill.background()
        shape.line.color.rgb = _hex_color("#7C3AED")
        shape.line.width = Pt(1.5)

    @staticmethod
    def _add_debug_region(
        slide: Any,
        region: dict[str, Any],
        block: dict[str, Any],
        index: int,
        font_family: str,
    ) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            _inches(region["x"], LOGICAL_WIDTH, SLIDE_WIDTH_INCHES),
            _inches(region["y"], LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES),
            _inches(region["w"], LOGICAL_WIDTH, SLIDE_WIDTH_INCHES),
            _inches(region["h"], LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES),
        )
        palette = ["#DBEAFE", "#DCFCE7", "#FEF3C7", "#FCE7F3"]
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_color(palette[(index - 1) % len(palette)])
        shape.line.color.rgb = _hex_color("#2563EB")
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.12)
        frame.margin_right = Inches(0.12)
        label = frame.paragraphs[0]
        label.text = f"{region['region_id']} → {region['block_id']}"
        label.font.size = Pt(10)
        label.font.bold = True
        label.font.color.rgb = _hex_color("#1D4ED8")
        for run in label.runs:
            _set_run_font(
                run,
                {
                    "font_family": font_family,
                    "font_size": 10,
                    "font_weight": 700,
                    "color": "#1D4ED8",
                },
            )
        content = frame.add_paragraph()
        content.text = " / ".join(_block_text(block["content"]))[:240]
        content.font.size = Pt(13)
        content.font.color.rgb = _hex_color("#111827")
        for run in content.runs:
            _set_run_font(
                run,
                {
                    "font_family": font_family,
                    "font_size": 13,
                    "font_weight": 400,
                    "color": "#111827",
                },
            )

    @staticmethod
    def _add_debug_label(slide: Any, plan: dict[str, Any]) -> None:
        box = slide.shapes.add_textbox(Inches(0.18), Inches(0.08), Inches(5.2), Inches(0.25))
        box.text_frame.text = (
            f"LAYOUT DEBUG · {plan['slide_id']} · {plan['layout_family']} · 1280×720"
        )
        box.text_frame.paragraphs[0].font.size = Pt(9)
        box.text_frame.paragraphs[0].font.color.rgb = _hex_color("#6B7280")

    @staticmethod
    def validate_output(output_path: Path, layouts: dict[str, Any]) -> None:
        presentation = Presentation(output_path)
        if len(presentation.slides) != len(layouts["plans"]):
            raise ValueError("Debug PPTX page count does not match Layout Plans")
        for slide, plan in zip(presentation.slides, layouts["plans"], strict=True):
            text = "\n".join(
                shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            )
            for region in plan["regions"]:
                if region["region_id"] not in text or region["block_id"] not in text:
                    raise ValueError(f"Debug PPTX lost mapping for {region['region_id']}")


class MinimalDesignPptxRenderBackend:
    """Render a minimal visual design as native editable PPTX text and shapes."""

    name = "python-pptx-minimal-design"
    version = "0.2.0"

    def render(self, request: RenderRequest) -> RenderResult:
        if request.target_format != "pptx":
            raise ValueError("MinimalDesignPptxRenderBackend only supports pptx")
        workspace = request.workspace.resolve()
        output_dir = ensure_within(workspace, request.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        specs = read_json(workspace / "slides/slide_specs.json")
        layouts = read_json(workspace / "layout/layout_plans.json")
        visual = read_json(workspace / "design/visual_system.json")
        outline = read_json(workspace / "outline/deck_outline.json")
        project_id = specs["project_id"]
        output_path = output_dir / f"{project_id.lower()}-final.pptx"

        presentation, blank_layout = _new_presentation()
        specs_by_id = {slide["slide_id"]: slide for slide in specs["slides"]}

        for plan in layouts["plans"]:
            slide_spec = specs_by_id[plan["slide_id"]]
            slide = presentation.slides.add_slide(blank_layout)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = _hex_color(visual["colors"]["background"])
            self._add_family_design(slide, plan, visual)
            blocks = {item["block_id"]: item for item in slide_spec["content_blocks"]}
            for region in sorted(plan["regions"], key=lambda item: item["z"]):
                self._add_region(slide, region, blocks[region["block_id"]], visual)
            self._add_page_number(slide, plan["slide_id"], visual)

        buffer = BytesIO()
        presentation.save(buffer)
        atomic_write_bytes(output_path, buffer.getvalue())
        self.validate_output(output_path, outline, specs)
        model_previews = self._render_model_previews(
            output_dir / "design-previews", specs, layouts, visual
        )
        return RenderResult(
            status="success",
            output_paths=(output_path,),
            preview_paths=tuple(model_previews),
            actual_editability_level="E3",
            warnings=(
                "Design SVGs are same-model proofs; independent Office previews remain required for G8.",
                "Minimal DesignImpl supports native text, panels, markers, and simple geometric shapes.",
            ),
        )

    @staticmethod
    def _add_family_design(
        slide: Any, plan: dict[str, Any], visual: dict[str, Any]
    ) -> None:
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(0.12),
            Inches(SLIDE_HEIGHT_INCHES),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = _hex_color(visual["colors"]["accent"])
        accent.line.fill.background()
        family = plan["layout_family"]
        if family == "hero":
            orb = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(10.7), Inches(-0.55), Inches(3.2), Inches(3.2)
            )
            orb.fill.solid()
            orb.fill.fore_color.rgb = _hex_color(visual["colors"]["primary"])
            orb.line.fill.background()
        elif family == "split":
            panel = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.75),
                Inches(1.85),
                Inches(4.45),
                Inches(4.15),
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = _hex_color(visual["colors"]["primary"])
            panel.line.fill.background()
            number = panel.text_frame.paragraphs[0]
            number.text = plan["slide_id"].split("-")[-1]
            number.font.size = Pt(70)
            number.font.bold = True
            number.font.color.rgb = _hex_color("#FFFFFF")
            number.alignment = PP_ALIGN.CENTER
            panel.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif family == "case":
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.72), Inches(2.05), Inches(0.8), Inches(0.8)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = _hex_color(visual["colors"]["accent"])
            marker.line.fill.background()
            paragraph = marker.text_frame.paragraphs[0]
            paragraph.text = plan["slide_id"].split("-")[-1]
            paragraph.font.size = Pt(17)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _hex_color("#FFFFFF")
            paragraph.alignment = PP_ALIGN.CENTER
            marker.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    @staticmethod
    def _add_page_number(slide: Any, slide_id: str, visual: dict[str, Any]) -> None:
        box = slide.shapes.add_textbox(Inches(12.25), Inches(7.05), Inches(0.65), Inches(0.2))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = slide_id.split("-")[-1]
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = _hex_color(visual["colors"]["text_secondary"])

    @staticmethod
    def _add_region(
        slide: Any,
        region: dict[str, Any],
        block: dict[str, Any],
        visual: dict[str, Any],
    ) -> None:
        x = _inches(region["x"], LOGICAL_WIDTH, SLIDE_WIDTH_INCHES)
        y = _inches(region["y"], LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES)
        width = _inches(region["w"], LOGICAL_WIDTH, SLIDE_WIDTH_INCHES)
        height = _inches(region["h"], LOGICAL_HEIGHT, SLIDE_HEIGHT_INCHES)
        is_surface = block["semantic_role"] in {"body", "evidence", "diagram", "table"}
        if is_surface:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_color(visual["colors"]["surface"])
            shape.line.color.rgb = _hex_color("#D8D2C6")
        else:
            shape = slide.shapes.add_textbox(x, y, width, height)
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.16)
        text_frame.margin_right = Inches(0.16)
        text_frame.margin_top = Inches(0.08)
        text_frame.margin_bottom = Inches(0.08)
        text_frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[region["valign"]]
        alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[region["align"]]

        style = _style_for(block, visual)
        values = _block_text(block["content"])
        for index, value in enumerate(values):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = value
            paragraph.alignment = alignment
            paragraph.level = 0
            paragraph.space_after = Pt(8 if len(values) > 1 else 0)
            if len(values) > 1 and block["content_type"] == "list":
                paragraph.text = f"• {value}"
            for run in paragraph.runs:
                _set_run_font(run, style)

    @staticmethod
    def validate_output(
        output_path: Path,
        outline: dict[str, Any],
        specs: dict[str, Any],
    ) -> None:
        """Reopen the generated file and verify native slide/text coverage."""

        presentation = Presentation(output_path)
        expected_slides = len(outline["slides"])
        if len(presentation.slides) != expected_slides:
            raise ValueError(
                f"PPTX slide count mismatch: expected {expected_slides}, got {len(presentation.slides)}"
            )
        for index, (slide, spec) in enumerate(
            zip(presentation.slides, specs["slides"], strict=True), start=1
        ):
            rendered_text = "\n".join(
                shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            )
            for block in spec["content_blocks"]:
                for value in _block_text(block["content"]):
                    if value not in rendered_text:
                        raise ValueError(
                            f"PPTX slide {index} is missing native text from {block['block_id']}"
                        )

    @staticmethod
    def _render_model_previews(
        output_dir: Path,
        specs: dict[str, Any],
        layouts: dict[str, Any],
        visual: dict[str, Any],
    ) -> list[Path]:
        output_dir.mkdir(parents=True, exist_ok=True)
        specs_by_id = {slide["slide_id"]: slide for slide in specs["slides"]}
        outputs: list[Path] = []
        for plan in layouts["plans"]:
            blocks = {
                block["block_id"]: block
                for block in specs_by_id[plan["slide_id"]]["content_blocks"]
            }
            svg = [
                '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">',
                f'<rect width="1280" height="720" fill="{_xml_escape(visual["colors"]["background"])}"/>',
                f'<rect width="12" height="720" fill="{_xml_escape(visual["colors"]["accent"])}"/>',
            ]
            family = plan["layout_family"]
            if family == "hero":
                svg.append(
                    f'<circle cx="1170" cy="65" r="155" fill="{_xml_escape(visual["colors"]["primary"])}"/>'
                )
            elif family == "split":
                svg.extend(
                    [
                        f'<rect x="72" y="178" width="428" height="398" rx="18" fill="{_xml_escape(visual["colors"]["primary"])}"/>',
                        f'<text x="286" y="420" text-anchor="middle" font-family="Arial,sans-serif" font-size="92" font-weight="700" fill="#FFFFFF">{_xml_escape(plan["slide_id"].split("-")[-1])}</text>',
                    ]
                )
            elif family == "case":
                svg.extend(
                    [
                        f'<circle cx="110" cy="245" r="40" fill="{_xml_escape(visual["colors"]["accent"])}"/>',
                        f'<text x="110" y="255" text-anchor="middle" font-family="Arial,sans-serif" font-size="20" font-weight="700" fill="#FFFFFF">{_xml_escape(plan["slide_id"].split("-")[-1])}</text>',
                    ]
                )
            for region in plan["regions"]:
                block = blocks[region["block_id"]]
                if block["semantic_role"] in {"body", "evidence", "diagram", "table"}:
                    svg.append(
                        f'<rect x="{region["x"]}" y="{region["y"]}" width="{region["w"]}" height="{region["h"]}" rx="10" fill="{_xml_escape(visual["colors"]["surface"])}" stroke="#D8D2C6"/>'
                    )
                style = _style_for(block, visual)
                y = region["y"] + style["font_size"] + 12
                for line in _block_text(block["content"]):
                    svg.append(
                        f'<text x="{region["x"] + 16}" y="{y}" font-family="Arial,sans-serif" font-size="{style["font_size"]}" fill="{_xml_escape(style["color"])}">{_xml_escape(line[:110])}</text>'
                    )
                    y += style["font_size"] * 1.35
            svg.append("</svg>")
            path = output_dir / f"{plan['slide_id']}.svg"
            path.write_text("\n".join(svg) + "\n", encoding="utf-8")
            outputs.append(path)
        return outputs


class MinimalPptxRenderBackend(MinimalDesignPptxRenderBackend):
    """Backward-compatible alias for the minimal final DesignImpl backend."""


class LibreOfficeDocumentRenderer:
    """Render PPTX through LibreOffice and Poppler for independent PNG previews."""

    def __init__(
        self,
        *,
        soffice: str | None = None,
        pdftoppm: str | None = None,
        font_match: str | None = None,
        timeout_seconds: int = 120,
    ) -> None:
        self.soffice = soffice or shutil.which("soffice")
        self.pdftoppm = pdftoppm or shutil.which("pdftoppm")
        self.font_match = font_match or shutil.which("fc-match")
        self.timeout_seconds = timeout_seconds

    @property
    def available(self) -> bool:
        return bool(self.soffice and self.pdftoppm)

    def preview(self, document_path: Path, output_dir: Path) -> Sequence[Path]:
        if not self.available:
            raise RuntimeError("LibreOffice preview requires soffice and pdftoppm")
        output_dir.mkdir(parents=True, exist_ok=True)
        with tempfile.TemporaryDirectory(prefix="slidethus-preview-") as temporary_name:
            temporary_dir = Path(temporary_name)
            profile_dir = temporary_dir / "libreoffice-profile"
            profile_dir.mkdir()
            cache_dir = temporary_dir / "font-cache"
            cache_dir.mkdir()
            system_font_paths = [
                path
                for path in (
                    "/System/Library/Fonts",
                    "/System/Library/Fonts/Supplemental",
                    "/Library/Fonts",
                    "/usr/share/fonts",
                    "/usr/local/share/fonts",
                )
                if Path(path).exists()
            ]
            process_environment = {
                **os.environ,
                "XDG_CACHE_HOME": str(cache_dir),
                "SAL_FONTPATH": os.pathsep.join(system_font_paths),
            }
            staged_fonts = self._stage_document_fonts(
                document_path,
                profile_dir,
                process_environment,
            )
            if self._contains_cjk(document_path) and staged_fonts == 0:
                raise RuntimeError(
                    "CJK preview requires a discoverable local font; no document font was staged"
                )
            convert = subprocess.run(
                [
                    str(self.soffice),
                    f"-env:UserInstallation={profile_dir.as_uri()}",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(temporary_dir),
                    str(document_path),
                ],
                check=False,
                capture_output=True,
                text=True,
                timeout=self.timeout_seconds,
                env=process_environment,
            )
            pdf_path = temporary_dir / f"{document_path.stem}.pdf"
            if convert.returncode != 0 or not pdf_path.exists():
                detail = (convert.stderr or convert.stdout).strip()
                raise RuntimeError(f"LibreOffice preview failed: {detail or convert.returncode}")
            prefix = output_dir / "slide"
            raster = subprocess.run(
                [str(self.pdftoppm), "-png", "-r", "144", str(pdf_path), str(prefix)],
                check=False,
                capture_output=True,
                text=True,
                timeout=self.timeout_seconds,
                env=process_environment,
            )
            if raster.returncode != 0:
                detail = (raster.stderr or raster.stdout).strip()
                raise RuntimeError(f"PDF preview rasterization failed: {detail or raster.returncode}")
        paths = sorted(output_dir.glob("slide-*.png"))
        if not paths:
            raise RuntimeError("Independent preview produced no PNG pages")
        for path in paths:
            if path.stat().st_size == 0:
                raise RuntimeError(f"Independent preview is empty: {path}")
        return tuple(paths)

    def _stage_document_fonts(
        self,
        document_path: Path,
        profile_dir: Path,
        environment: dict[str, str],
    ) -> int:
        if not self.font_match:
            return 0
        families: set[str] = set()
        presentation = Presentation(document_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            families.add(run.font.name)
        font_dir = profile_dir / "user" / "fonts"
        font_dir.mkdir(parents=True, exist_ok=True)
        staged: set[Path] = set()
        for family in sorted(families):
            match = subprocess.run(
                [self.font_match, "-f", "%{file}\n", family],
                check=False,
                capture_output=True,
                text=True,
                timeout=10,
                env=environment,
            )
            if match.returncode != 0:
                continue
            first_line = match.stdout.splitlines()[0].strip() if match.stdout else ""
            source = Path(first_line)
            if not source.is_file() or source.suffix.lower() not in {".ttf", ".otf", ".ttc"}:
                continue
            destination = font_dir / source.name
            if destination not in staged:
                shutil.copyfile(source, destination)
                staged.add(destination)
        return len(staged)

    @staticmethod
    def _contains_cjk(document_path: Path) -> bool:
        presentation = Presentation(document_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and re.search(
                    r"[\u3400-\u9fff]", shape.text
                ):
                    return True
        return False
