from __future__ import annotations

import json
import os
import shutil
from pathlib import Path
from typing import Any

import pytest
from pptx import Presentation

from slidethus.artifact_runtime import ArtifactRuntime
from slidethus.cli import main
from slidethus.protocols import BriefCompletionHints
from slidethus.semantic_reviews import SEMANTIC_DIMENSIONS
from slidethus.services.brief_completion import BriefCompletionService
from slidethus.services.workflow_application import (
    WorkflowApplicationService,
    WorkflowRequest,
)
from slidethus.state_machine import Phase
from slidethus.validation import validate_workspace
from slidethus.workflow_application_reports import workflow_request_hash
from slidethus.workflow_operations import (
    WorkflowLease,
    WorkflowOperationalLimits,
    persist_workflow_event,
    utc_now,
    workflow_attempt_id,
)
from slidethus.workspace import init_workspace
from tests.fontconfig_fakes import write_fontconfig_tools


class CleanSemanticProvider:
    name = "workflow-clean-semantic"
    version = "1.0.0"

    def review(self, context: dict[str, Any]) -> dict[str, Any]:
        if context["mode"] == "open_issue":
            return {"issues": []}
        return {
            "dimensions": [
                {
                    "dimension": dimension,
                    "score": 5,
                    "rationale": "Workflow fixture quality baseline.",
                    "issue_ids": [],
                }
                for dimension in SEMANTIC_DIMENSIONS
            ]
        }


class CleanVisualProvider:
    name = "workflow-clean-visual"
    version = "1.0.0"

    def review(self, image_paths: tuple[Path, ...], context: dict[str, Any]) -> dict[str, Any]:
        assert image_paths
        return {"issues": []}


class RisingCostMeter:
    name = "rising-cost-meter"
    version = "1.0.0"

    def __init__(self) -> None:
        self.calls = 0

    def current_cost_usd(self) -> float:
        self.calls += 1
        return 0.0 if self.calls == 1 else 2.0


def _renderer_root() -> Path | None:
    raw = os.environ.get("SLIDETHUS_PPTXGENJS_TEST_ROOT")
    return Path(raw).resolve() if raw else None


def _font_match(tmp_path: Path) -> Path:
    return write_fontconfig_tools(tmp_path)


def _hints() -> BriefCompletionHints:
    return BriefCompletionHints(
        request_text="Create an 8-page management decision deck about an enterprise agent operating model",
        purpose="Present the enterprise agent operating model",
        desired_outcome="Approve implementation",
        call_to_action="Approve project initiation",
        delivery_context="Management decision meeting",
        audience_role="Executive management",
        page_target=8,
    )


def _service(workspace: Path, tmp_path: Path) -> WorkflowApplicationService:
    return WorkflowApplicationService(
        workspace,
        semantic_provider=CleanSemanticProvider(),
        visual_provider=CleanVisualProvider(),
        renderer_root=_renderer_root(),
        font_match=str(_font_match(tmp_path)),
    )


def _source(tmp_path: Path) -> Path:
    path = tmp_path / "source.md"
    path.write_text(
        "# Enterprise operating model\n\n"
        "Enterprises build data, knowledge, process, rules, tools, permissions and evaluation standards.\n\n"
        "# Risk\n\nAdding more agents does not automatically improve task quality.\n",
        encoding="utf-8",
    )
    return path


def _pptx_source(tmp_path: Path) -> Path:
    path = tmp_path / "existing-deck.pptx"
    presentation = Presentation()
    title = presentation.slides.add_slide(presentation.slide_layouts[1])
    title.shapes.title.text = "Enterprise agent operating model"
    title.placeholders[1].text = "Data, knowledge, process, rules, tools, permissions and evaluation standards"
    risk = presentation.slides.add_slide(presentation.slide_layouts[1])
    risk.shapes.title.text = "Risk"
    risk.placeholders[1].text = "Adding more agents does not automatically improve task quality"
    presentation.save(path)
    return path


@pytest.fixture(scope="module")
def created_baseline(tmp_path_factory: pytest.TempPathFactory) -> tuple[Path, dict[str, Any]]:
    if _renderer_root() is None:
        pytest.skip("real M4 sidecar is required for M6.1 workflow integration")
    base = tmp_path_factory.mktemp("workflow-create")
    workspace = base / "workspace"
    service = _service(workspace, base)
    result = service.run(
        WorkflowRequest(
            workflow="create",
            title="Workflow Create",
            source_paths=(_source(base),),
            brief_hints=_hints(),
        )
    )
    return workspace, result.report


def test_create_workflow_reaches_g8_and_persists_valid_report(created_baseline: tuple[Path, dict[str, Any]]) -> None:
    workspace, report = created_baseline

    assert report["status"] == "ready"
    assert report["workflow"] == "create"
    assert report["gate_result"] == {"gate_id": "G8", "status": "pass", "reasons": []}
    assert {item["kind"] for item in report["outputs"]} == {
        "m3_application",
        "m4_application",
        "m5_application",
    }
    assert validate_workspace(workspace, check_hashes=True).ok


def test_audit_workflow_preserves_frozen_truth(created_baseline: tuple[Path, dict[str, Any]], tmp_path: Path) -> None:
    baseline, _ = created_baseline
    workspace = tmp_path / "workspace"
    shutil.copytree(baseline, workspace)
    before = json.loads((workspace / "project_state.json").read_text(encoding="utf-8"))
    frozen_before = {
        item["artifact_type"]: (item["version"], item["content_hash"])
        for item in before["artifacts"]
        if item["artifact_type"] not in {"quality_report", "gate_results", "decision_log", "assumption_log"}
    }

    result = _service(workspace, tmp_path).run(
        WorkflowRequest(workflow="audit", title="Workflow Audit", auto_repair=False)
    )
    after = json.loads((workspace / "project_state.json").read_text(encoding="utf-8"))
    frozen_after = {
        item["artifact_type"]: (item["version"], item["content_hash"])
        for item in after["artifacts"]
        if item["artifact_type"] not in {"quality_report", "gate_results", "decision_log", "assumption_log"}
    }

    assert result.report["status"] == "ready"
    assert result.report["changed_artifacts"] == []
    assert frozen_after == frozen_before
    assert validate_workspace(workspace, check_hashes=True).ok


def test_improve_workflow_audits_before_repair_and_is_noop_on_clean_deck(created_baseline: tuple[Path, dict[str, Any]], tmp_path: Path) -> None:
    baseline, _ = created_baseline
    workspace = tmp_path / "workspace"
    shutil.copytree(baseline, workspace)

    result = _service(workspace, tmp_path).run(
        WorkflowRequest(workflow="improve", title="Workflow Improve")
    )

    assert result.report["status"] == "ready"
    assert result.report["changed_artifacts"] == []
    assert any(item["stage"] == "improve_audit" for item in result.report["actions"])
    assert validate_workspace(workspace, check_hashes=True).ok


@pytest.mark.skipif(
    _renderer_root() is None,
    reason="real M4 sidecar is required for M6.1 rebuild integration",
)
def test_rebuild_workflow_keeps_original_pptx_byte_identical(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    before = source.read_bytes()
    workspace = tmp_path / "rebuild-workspace"

    result = _service(workspace, tmp_path).run(
        WorkflowRequest(
            workflow="rebuild",
            title="Workflow Rebuild",
            source_paths=(source,),
            brief_hints=_hints(),
        )
    )

    assert result.report["status"] == "ready"
    assert result.report["workflow"] == "rebuild"
    assert source.read_bytes() == before
    assert any(item["stage"] == "source_preservation" for item in result.report["actions"])
    assert validate_workspace(workspace, check_hashes=True).ok


def test_extract_style_produces_schema_valid_candidate_without_copying_reference_bytes(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    before = source.read_bytes()
    workspace = tmp_path / "style-workspace"

    result = _service(workspace, tmp_path).run(
        WorkflowRequest(
            workflow="extract_style",
            title="Workflow Style",
            source_paths=(source,),
        )
    )

    assert result.report["status"] == "ready"
    assert result.report["changed_artifacts"] == []
    output = next(item for item in result.report["outputs"] if item["kind"] == "visual_system_candidate")
    candidate = json.loads((workspace / output["path"]).read_text(encoding="utf-8"))
    assert candidate["theme_id"].startswith("THEME-EXTRACTED-")
    assert candidate["forbidden_patterns"] == [
        "copy-unlicensed-font-bytes",
        "copy-unlicensed-brand-assets",
        "pixel-template-lock",
    ]
    assert source.read_bytes() == before
    assert validate_workspace(workspace, check_hashes=True).ok


def test_extract_style_cli_run_list_show(tmp_path: Path, capsys) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "style-cli-workspace"

    assert main([
        "workflow", "run", "extract_style", str(workspace),
        "--source", str(source),
        "--title", "Style CLI",
    ]) == 0
    output = json.loads(capsys.readouterr().out)
    report_id = output["report_id"]
    assert output["report"]["status"] == "ready"
    assert main(["workflow", "list", str(workspace)]) == 0
    assert report_id in capsys.readouterr().out
    assert main(["workflow", "show", str(workspace), report_id]) == 0
    assert '"workflow": "extract_style"' in capsys.readouterr().out


def test_workflow_cache_records_miss_then_hit(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "style-cache-workspace"
    service = _service(workspace, tmp_path)
    request = WorkflowRequest(
        workflow="extract_style",
        title="Style Cache",
        source_paths=(source,),
    )

    first = service.run(request)
    second = service.run(request)
    operations = [
        json.loads(path.read_text(encoding="utf-8"))
        for path in (workspace / ".slidethus/workflows/operations").glob("*.json")
    ]

    assert first.report == second.report
    assert first.path == second.path
    assert second.changed is False
    assert {item["cache_status"] for item in operations} == {"miss", "hit"}
    assert validate_workspace(workspace, check_hashes=True).ok


def test_workflow_cache_can_be_disabled_by_zero_ttl(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "style-no-cache-workspace"
    service = WorkflowApplicationService(
        workspace,
        operational_limits=WorkflowOperationalLimits(max_cache_age_seconds=0),
    )
    request = WorkflowRequest(
        workflow="extract_style",
        title="Style No Cache",
        source_paths=(source,),
    )

    first = service.run(request)
    second = service.run(request)
    operations = [
        json.loads(path.read_text(encoding="utf-8"))
        for path in (workspace / ".slidethus/workflows/operations").glob("*.json")
    ]

    assert first.report == second.report
    assert all(item["cache_status"] == "miss" for item in operations)
    assert len(operations) == 2
    assert validate_workspace(workspace, check_hashes=True).ok


def test_workflow_cache_invalidates_when_bound_workspace_artifact_changes(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "style-currentness-workspace"
    service = WorkflowApplicationService(workspace)
    request = WorkflowRequest(
        workflow="extract_style",
        title="Style Currentness",
        source_paths=(source,),
    )

    service.run(request)
    runtime = ArtifactRuntime(workspace)
    brief, version = runtime.read_artifact_snapshot("project_brief")
    brief["intent"]["purpose"] = "Invalidate the previous workflow cache binding"
    runtime.write_artifact(
        "project_brief",
        brief,
        expected_version=version,
        created_by="workflow-cache-currentness-test",
    )
    second = service.run(request)
    operations = [
        json.loads(path.read_text(encoding="utf-8"))
        for path in (workspace / ".slidethus/workflows/operations").glob("*.json")
    ]

    assert second.report["status"] == "ready"
    assert all(item["cache_status"] == "miss" for item in operations)
    assert len(operations) == 2
    assert validate_workspace(workspace, check_hashes=True).ok


def test_orphan_workflow_attempt_is_recovered_before_next_run(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = init_workspace(tmp_path / "recovery-workspace", title="Recovery Workspace")
    service = WorkflowApplicationService(workspace)
    request = WorkflowRequest(
        workflow="extract_style",
        title="Recovery Style",
        source_paths=(source,),
    )
    request_hash = workflow_request_hash(service._request_payload(request))
    execution_signature = service._execution_signature()
    project_id = json.loads((workspace / "project_state.json").read_text(encoding="utf-8"))["project_id"]
    started_at = utc_now()
    attempt_id = workflow_attempt_id(
        project_id,
        "extract_style",
        request_hash,
        execution_signature,
        started_at,
    )
    persist_workflow_event(
        workspace,
        schema_dir=service.schemas.schema_dir,
        project_id=project_id,
        attempt_id=attempt_id,
        workflow="extract_style",
        request_hash=request_hash,
        execution_signature=execution_signature,
        event_type="started",
        sequence=1,
        detail="Fixture simulates a process interruption after lease admission.",
    )

    before = validate_workspace(workspace, check_hashes=True)
    assert before.ok

    result = service.run(request)
    events = [
        json.loads(path.read_text(encoding="utf-8"))
        for path in (workspace / ".slidethus/workflows/events").glob("*.json")
    ]
    recovered = [item for item in events if item["attempt_id"] == attempt_id]

    assert result.report["status"] == "ready"
    assert {item["event_type"] for item in recovered} == {"started", "recovered"}
    assert validate_workspace(workspace, check_hashes=True).ok


def test_create_resume_admission_is_explicit(tmp_path: Path) -> None:
    workspace = init_workspace(tmp_path / "resume-admission-workspace", title="Resume Admission")
    completed = BriefCompletionService(workspace).complete(_hints())
    assert completed.status == "resolved"
    ArtifactRuntime(workspace).record_gate("G0", target_phase=Phase.BRIEF_READY)
    service = WorkflowApplicationService(workspace)
    request = WorkflowRequest(
        workflow="create",
        title="Workflow Create",
        source_paths=(_source(tmp_path),),
        brief_hints=_hints(),
    )

    with pytest.raises(Exception, match="requires a new/stage-0 workspace"):
        service._ensure_workspace(request)
    service._ensure_workspace(request, allow_resume=True)


def test_workflow_wall_time_overrun_returns_blocked_operation(tmp_path: Path) -> None:
    class ExpiredDeadlineService(WorkflowApplicationService):
        def _run_once(
            self,
            request: WorkflowRequest,
            *,
            allow_resume: bool = False,
        ):
            del request, allow_resume
            self._deadline_ns = 0
            self._check_deadline("fixture")
            raise AssertionError("deadline check must block")

    workspace = init_workspace(tmp_path / "wall-budget-workspace", title="Wall Budget Workspace")
    service = ExpiredDeadlineService(
        workspace,
        operational_limits=WorkflowOperationalLimits(max_wall_seconds=1),
    )

    result = service.run(WorkflowRequest(workflow="audit", title="Wall Budget Audit"))
    operation = next(
        json.loads(path.read_text(encoding="utf-8"))
        for path in (workspace / ".slidethus/workflows/operations").glob("*.json")
    )

    assert result.report["status"] == "blocked"
    assert result.report["blockers"][0]["code"] == "workflow_wall_time_budget_exceeded"
    assert operation["status"] == "blocked"
    assert validate_workspace(workspace, check_hashes=True).ok


def test_workflow_input_budget_blocks_before_style_extraction(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "style-budget-workspace"
    service = WorkflowApplicationService(
        workspace,
        operational_limits=WorkflowOperationalLimits(
            max_input_bytes=1,
            max_slide_updates=64,
            max_wall_seconds=900,
        ),
    )

    result = service.run(
        WorkflowRequest(
            workflow="extract_style",
            title="Style Budget",
            source_paths=(source,),
        )
    )

    assert result.report["status"] == "blocked"
    assert result.report["blockers"][0]["code"] == "workflow_input_budget_exceeded"
    assert not (workspace / ".slidethus/workflows/style-candidates").exists()
    assert validate_workspace(workspace, check_hashes=True).ok


def test_structured_request_bytes_are_included_in_input_budget(tmp_path: Path) -> None:
    workspace = init_workspace(tmp_path / "request-budget-workspace", title="Request Budget Workspace")
    service = WorkflowApplicationService(
        workspace,
        operational_limits=WorkflowOperationalLimits(max_input_bytes=512),
    )

    result = service.run(
        WorkflowRequest(
            workflow="revise",
            title="Request Budget",
            slide_updates={"S-001": {"headline": "X" * 5000}},
            reason="Budget admission must count structured request bytes",
        )
    )

    assert result.report["status"] == "blocked"
    assert result.report["blockers"][0]["code"] == "workflow_input_budget_exceeded"
    assert not (workspace / "outline/deck_outline.json").exists()
    assert validate_workspace(workspace, check_hashes=True).ok


def test_provider_cost_budget_blocks_when_meter_is_unavailable(tmp_path: Path) -> None:
    workspace = init_workspace(tmp_path / "cost-workspace", title="Cost Workspace")
    service = WorkflowApplicationService(
        workspace,
        semantic_provider=CleanSemanticProvider(),
        visual_provider=CleanVisualProvider(),
        operational_limits=WorkflowOperationalLimits(
            max_input_bytes=100 * 1024 * 1024,
            max_slide_updates=64,
            max_wall_seconds=900,
            max_provider_cost_usd=1.0,
        ),
    )

    result = service.run(WorkflowRequest(workflow="audit", title="Cost Audit"))

    assert result.report["status"] == "blocked"
    assert result.report["blockers"] == [
        {
            "code": "workflow_cost_meter_missing",
            "message": "A provider-cost budget was requested, but no WorkflowCostMeter was injected.",
        }
    ]
    assert validate_workspace(workspace, check_hashes=True).ok


def test_provider_cost_overrun_returns_blocked_operation(tmp_path: Path) -> None:
    workspace = init_workspace(tmp_path / "measured-cost-workspace", title="Measured Cost Workspace")
    meter = RisingCostMeter()
    service = WorkflowApplicationService(
        workspace,
        semantic_provider=CleanSemanticProvider(),
        cost_meter=meter,
        operational_limits=WorkflowOperationalLimits(max_provider_cost_usd=1.0),
    )

    result = service.run(WorkflowRequest(workflow="audit", title="Measured Cost Audit"))
    operation = next(
        json.loads(path.read_text(encoding="utf-8"))
        for path in (workspace / ".slidethus/workflows/operations").glob("*.json")
    )

    assert result.report["status"] == "blocked"
    assert result.report["blockers"][0]["code"] == "workflow_provider_cost_budget_exceeded"
    assert operation["status"] == "blocked"
    assert operation["metrics"]["provider_cost_status"] == "measured"
    assert operation["metrics"]["provider_cost_usd"] >= 2.0
    assert validate_workspace(workspace, check_hashes=True).ok


def test_workflow_exclusive_lease_rejects_concurrent_mutator(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = init_workspace(tmp_path / "lease-workspace", title="Lease Workspace")
    service = _service(workspace, tmp_path)

    with WorkflowLease(workspace):
        with pytest.raises(Exception, match="already holds the workspace lease"):
            service.run(
                WorkflowRequest(
                    workflow="extract_style",
                    title="Lease Style",
                    source_paths=(source,),
                )
            )


def test_workflow_operation_tampering_is_detected(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "operation-tamper-workspace"
    _service(workspace, tmp_path).run(
        WorkflowRequest(
            workflow="extract_style",
            title="Operation Tamper",
            source_paths=(source,),
        )
    )
    operation = next((workspace / ".slidethus/workflows/operations").glob("*.json"))
    data = json.loads(operation.read_text(encoding="utf-8"))
    data["duration_ms"] += 1
    operation.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    validation = validate_workspace(workspace, check_hashes=True)
    assert not validation.ok
    assert any(item.code == "invalid_workflow_operation_report" for item in validation.issues)
    with pytest.raises(Exception, match="operational history is invalid"):
        _service(workspace, tmp_path).run(
            WorkflowRequest(
                workflow="extract_style",
                title="Operation Tamper",
                source_paths=(source,),
            )
        )


def test_workflow_event_tampering_is_detected(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "event-tamper-workspace"
    _service(workspace, tmp_path).run(
        WorkflowRequest(
            workflow="extract_style",
            title="Event Tamper",
            source_paths=(source,),
        )
    )
    event = next((workspace / ".slidethus/workflows/events").glob("*.json"))
    data = json.loads(event.read_text(encoding="utf-8"))
    data["detail"] = "tampered event detail"
    event.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    validation = validate_workspace(workspace, check_hashes=True)
    assert not validation.ok
    assert any(item.code == "invalid_workflow_event" for item in validation.issues)


def test_workflow_report_tampering_is_detected(tmp_path: Path) -> None:
    source = _pptx_source(tmp_path)
    workspace = tmp_path / "style-tamper-workspace"
    result = _service(workspace, tmp_path).run(
        WorkflowRequest(
            workflow="extract_style",
            title="Style Tamper",
            source_paths=(source,),
        )
    )
    data = json.loads(result.path.read_text(encoding="utf-8"))
    data["final_phase"] = "COMPLETED"
    result.path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    validation = validate_workspace(workspace, check_hashes=True)
    assert not validation.ok
    assert any(item.code == "invalid_workflow_application_report" for item in validation.issues)


@pytest.mark.skipif(
    _renderer_root() is None,
    reason="real M4 sidecar is required for M6.1 revise integration",
)
def test_revise_workflow_updates_target_slide_and_regenerates_downstream(created_baseline: tuple[Path, dict[str, Any]], tmp_path: Path) -> None:
    baseline, _ = created_baseline
    workspace = tmp_path / "workspace"
    shutil.copytree(baseline, workspace)
    outline = json.loads((workspace / "outline/deck_outline.json").read_text(encoding="utf-8"))
    target = str(outline["slides"][0]["slide_id"])
    headline = "Revised enterprise agent operating model"

    result = _service(workspace, tmp_path).run(
        WorkflowRequest(
            workflow="revise",
            title="Workflow Revise",
            slide_updates={target: {"headline": headline}},
            reason="Clarify the opening proposition",
        )
    )

    current = json.loads((workspace / "outline/deck_outline.json").read_text(encoding="utf-8"))
    revised = next(item for item in current["slides"] if item["slide_id"] == target)
    assert result.report["status"] == "ready"
    assert revised["headline"] == headline
    assert "deck_outline" in result.report["changed_artifacts"]
    assert "render_manifest" in result.report["changed_artifacts"]
    assert any(item["stage"] == "revise_regression" for item in result.report["actions"])
    assert validate_workspace(workspace, check_hashes=True).ok
